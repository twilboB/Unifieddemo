"""
Inspect PowerPoint file structure using python-pptx.
Prints detailed info about each slide, shapes, placeholders, and charts.
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

PPTX_PATH = "/Users/timwilson/Downloads/PCA_Coles_Black Friday _ Cyber Monday  -  Repaired.pptx"

def shape_type_str(shape):
    """Return a readable string for shape type."""
    try:
        return str(shape.shape_type)
    except Exception:
        return "UNKNOWN"

def get_text(shape):
    """Get full text from a shape, handling nested frames."""
    try:
        if shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in para.runs)
                if para_text:
                    lines.append(para_text)
            return "\n".join(lines) if lines else shape.text_frame.text
        return ""
    except Exception as e:
        return f"[TEXT ERROR: {e}]"

def pos_str(shape):
    """Return position/size as readable string."""
    try:
        return (f"left={shape.left}, top={shape.top}, "
                f"width={shape.width}, height={shape.height}")
    except Exception:
        return "N/A"

def inspect_shape(shape, indent="  "):
    """Print details of a shape."""
    name = shape.name
    stype = shape_type_str(shape)
    is_ph = shape.is_placeholder
    ph_idx = None
    ph_type = None
    if is_ph:
        try:
            ph_idx = shape.placeholder_format.idx
            ph_type = str(shape.placeholder_format.type)
        except Exception:
            ph_idx = "?"
            ph_type = "?"

    text = get_text(shape)
    text_preview = text[:200].replace("\n", " | ") if text else ""

    print(f"{indent}Shape: '{name}'")
    print(f"{indent}  type={stype}, is_placeholder={is_ph}", end="")
    if is_ph:
        print(f", ph_idx={ph_idx}, ph_type={ph_type}", end="")
    print()
    print(f"{indent}  pos: {pos_str(shape)}")
    if text_preview:
        print(f"{indent}  text: '{text_preview}'")
    else:
        print(f"{indent}  text: (empty)")

    # Charts
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        try:
            chart = shape.chart
            print(f"{indent}  CHART: chart_type={chart.chart_type}, "
                  f"series_count={len(chart.series)}")
            for i, series in enumerate(chart.series):
                print(f"{indent}    Series {i}: name={series.name}")
        except Exception as e:
            print(f"{indent}  CHART ERROR: {e}")

    # Groups — recurse
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for sub in shape.shapes:
                inspect_shape(sub, indent + "  ")
        except Exception as e:
            print(f"{indent}  GROUP ERROR: {e}")

def main():
    prs = Presentation(PPTX_PATH)
    print(f"=== Presentation: {os.path.basename(PPTX_PATH)} ===")
    print(f"Slides: {len(prs.slides)}, Slide size: "
          f"{prs.slide_width} x {prs.slide_height} EMU\n")

    for slide_idx, slide in enumerate(prs.slides):
        layout = slide.slide_layout
        layout_name = layout.name if layout else "N/A"
        print(f"{'='*70}")
        print(f"SLIDE {slide_idx}  |  Layout: '{layout_name}'")
        print(f"{'='*70}")

        shapes = slide.shapes
        print(f"  Shapes on slide: {len(shapes)}")

        # Detect potential duplicate titles: gather all shapes with text
        # at similar vertical positions
        title_candidates = []
        for shape in shapes:
            if shape.is_placeholder:
                try:
                    idx = shape.placeholder_format.idx
                    if idx == 0 or idx == 1:
                        title_candidates.append(shape)
                except Exception:
                    pass

        for shape in shapes:
            inspect_shape(shape)

        # Duplicate title detection
        # Look for multiple shapes with similar top positions and non-empty text
        positioned_shapes = []
        for shape in shapes:
            text = get_text(shape)
            if text.strip():
                try:
                    positioned_shapes.append((shape.top, shape.name, text[:80]))
                except Exception:
                    pass
        positioned_shapes.sort(key=lambda x: x[0])

        # Check for shapes with nearly identical top positions (within 50000 EMU ~ 0.5pt)
        dup_groups = []
        used = set()
        for i, (top_i, name_i, text_i) in enumerate(positioned_shapes):
            if i in used:
                continue
            group = [(top_i, name_i, text_i)]
            for j, (top_j, name_j, text_j) in enumerate(positioned_shapes):
                if j <= i or j in used:
                    continue
                if abs(top_i - top_j) < 200000:  # ~2pt overlap zone
                    group.append((top_j, name_j, text_j))
                    used.add(j)
            if len(group) > 1:
                dup_groups.append(group)
            used.add(i)

        if dup_groups:
            print(f"\n  *** POTENTIAL DUPLICATE/OVERLAPPING TEXT SHAPES ***")
            for grp in dup_groups:
                print(f"    Group at top~{grp[0][0]}:")
                for top, name, text in grp:
                    print(f"      '{name}' top={top}: '{text}'")

        # Also check layout shapes for inherited placeholders
        print(f"\n  Layout shapes (inherited, not directly on slide):")
        layout_ph_idxs = set()
        for lsh in layout.placeholders:
            layout_ph_idxs.add(lsh.placeholder_format.idx)

        slide_ph_idxs = set()
        for sh in slide.placeholders:
            slide_ph_idxs.add(sh.placeholder_format.idx)

        inherited = layout_ph_idxs - slide_ph_idxs
        if inherited:
            print(f"    Layout placeholder idx(s) NOT overridden on slide: {inherited}")
            for lsh in layout.placeholders:
                idx = lsh.placeholder_format.idx
                if idx in inherited:
                    lt = get_text(lsh)
                    print(f"      Layout ph idx={idx} name='{lsh.name}' text='{lt[:80]}'")
        else:
            print(f"    All layout placeholders are overridden on slide.")

        print()

    print("=== INSPECTION COMPLETE ===")

if __name__ == "__main__":
    main()
