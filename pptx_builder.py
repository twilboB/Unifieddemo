"""
Template-based PowerPoint generator for PCA Builder.
Copies and modifies OMD template.pptx rather than building from scratch.

Template slide map (0-based):
  0 → Cover        (Cover 1 layout)
  1 → Agenda       (Agenda 5 layout — 5 section slots)
  2 → Divider      (Divider layout — section break)
  3 → Chart slide  (1/3 Layout — left text, right chart)
  4 → Rec slide    (4 Col — what worked / didn't / tests / verdict)
  5 → Callout      (Big Text Color 1 — one bold statement)
"""
import copy, os, re, zipfile, tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "OMD template.pptx")

# Template slide indices
TMPL_COVER    = 0
TMPL_AGENDA   = 1
TMPL_DIVIDER  = 2
TMPL_CHART    = 3
TMPL_REC      = 4
TMPL_CALLOUT  = 5

# Chart bounding box from "1/3 Layout" slide (right 2/3 of slide)
CHART_L = Inches(4.665)
CHART_T = Inches(1.451)
CHART_W = Inches(8.33)
CHART_H = Inches(5.382)

CHART_TYPE_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
    "line":   XL_CHART_TYPE.LINE,
}

NO_CHART_SLIDE_TYPES = {"intro", "recommendation"}

# Relationship namespace
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# ─── Brand colours (OMD palette) ───────────────────────────────────────────────
BRAND_COLORS = [
    RGBColor(0xF0, 0x5C, 0x2C),  # orange
    RGBColor(0x43, 0x38, 0xCA),  # purple
    RGBColor(0x0C, 0x0A, 0x28),  # navy
    RGBColor(0x6D, 0x63, 0xE8),  # purple-lt
    RGBColor(0xFC, 0xA9, 0x81),  # orange-lt
    RGBColor(0x94, 0xA3, 0xB8),  # slate
]


def _axis_fmt(series_label: str) -> str:
    """Return an Excel number format string based on the series label."""
    sl = (series_label or "").lower()
    if "$" in sl or "spend" in sl or "cpm" in sl or "cpc" in sl or "cpcv" in sl:
        return '$#,##0.##'
    if "%" in sl or "ctr" in sl or "vtr" in sl:
        return '0.00%'
    if "impression" in sl or "click" in sl or "view" in sl:
        return '#,##0'
    return '#,##0.##'


def _apply_chart_style(chart, xl_type, series_label: str = ""):
    """
    Apply brand colours, axis number formatting, and line-chart cleanup
    (no markers, 1.75pt line) to a chart object.
    """
    # Value-axis number format
    try:
        chart.value_axis.number_format = _axis_fmt(series_label)
        chart.value_axis.number_format_is_linked = False
    except Exception:
        pass

    is_line = (xl_type == XL_CHART_TYPE.LINE)

    for i, series in enumerate(chart.series):
        clr = BRAND_COLORS[i % len(BRAND_COLORS)]
        try:
            if is_line:
                # Line colour + width
                series.format.line.color.rgb = clr
                series.format.line.width = Pt(1.75)
                # Suppress markers via XML
                ser_elem = series._element
                marker = ser_elem.find(qn('c:marker'))
                if marker is None:
                    marker = etree.SubElement(ser_elem, qn('c:marker'))
                symbol = marker.find(qn('c:symbol'))
                if symbol is None:
                    symbol = etree.SubElement(marker, qn('c:symbol'))
                symbol.set('val', 'none')
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = clr
        except Exception:
            pass

# ─── Placeholder index maps ────────────────────────────────────────────────────

# Agenda: 5 rows of (title_idx, number_idx)
AGENDA_ROWS = [(10, 11), (42, 43), (44, 45), (46, 47), (48, 49)]

# Divider
DIVIDER_NUM_IDX   = 16
DIVIDER_TITLE_IDX = 14
DIVIDER_SUB_IDX   = 15

# Rec (4 Col): (num_idx, body_idx)
# Note: simplified template has no per-column title placeholder
REC_COLUMNS = [
    (40, 47),
    (51, 53),
    (54, 56),
    (48, 50),
]

# Callout (Big Text Color 1) — only one text placeholder in simplified template
CALLOUT_STATEMENT_IDX = 42

# Layout placeholder indices to suppress on 1/3 Layout slides.
# ph=11  → "Chart goes here" instruction text in layout
# ph=14  → "Section Number/Title" instruction text in layout
# NOTE: ph=0 (narrow title) is NOT suppressed — the full-width Title 1 textbox
# sits directly on top of it, covering it in both edit and slide-show mode.
# NOTE: ph=172 (left panel) is intentionally excluded — the template slide's
# <p:bg> white override is now copied on clone, so ph=172 renders transparent.
_CHART_SLIDE_SUPPRESS = [11, 14]


# ─── Helpers ───────────────────────────────────────────────────────────────────

def _clean_values(raw):
    out = []
    for v in raw:
        try:
            out.append(float(str(v).replace(",", "").replace("$", "").replace("%", "")))
        except (ValueError, TypeError):
            out.append(0.0)
    return out


def _ph(slide, idx):
    """Find placeholder by idx on a slide."""
    for s in slide.shapes:
        if s.is_placeholder and s.placeholder_format.idx == idx:
            return s
    return None


def _shape(slide, name):
    """Find shape by name on a slide."""
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def _has_image_ref(elem):
    """
    Return True if elem (or any descendant) has an r:embed or r:link attribute.
    These point to relationship parts that would be invalid on a cloned slide.
    """
    embed_attr = f"{{{_R_NS}}}embed"
    link_attr  = f"{{{_R_NS}}}link"
    for e in elem.iter():
        if embed_attr in e.attrib or link_attr in e.attrib:
            return True
    return False


def _set_text(shape, text):
    """
    Set text on a shape, preserving the first run's formatting (rPr).
    Collapses the text frame to a single paragraph.
    """
    if shape is None or not shape.has_text_frame:
        return
    tf   = shape.text_frame
    txBody = tf._txBody
    first_p = txBody.find(qn('a:p'))
    if first_p is None:
        return

    # Grab the first run's rPr to preserve font/colour/size
    saved_rPr = None
    first_r = first_p.find(qn('a:r'))
    if first_r is not None:
        rPr = first_r.find(qn('a:rPr'))
        if rPr is not None:
            saved_rPr = copy.deepcopy(rPr)

    # Also preserve paragraph pPr (alignment, spacing)
    saved_pPr = None
    pPr = first_p.find(qn('a:pPr'))
    if pPr is not None:
        saved_pPr = copy.deepcopy(pPr)

    # Remove all paragraphs
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    # Rebuild first paragraph
    p_new = etree.SubElement(txBody, qn('a:p'))
    if saved_pPr is not None:
        p_new.append(saved_pPr)
    r_new = etree.SubElement(p_new, qn('a:r'))
    if saved_rPr is not None:
        r_new.append(saved_rPr)
    t_new = etree.SubElement(r_new, qn('a:t'))
    t_new.text = str(text)


def _set_bullets(shape, bullets):
    """
    Replace text frame with bullet list, preserving the template paragraph's
    formatting (pPr + rPr) as a template for all new paragraphs.
    """
    if not bullets or shape is None or not shape.has_text_frame:
        return
    tf     = shape.text_frame
    txBody = tf._txBody

    # Extract formatting templates from the first existing paragraph
    first_p = txBody.find(qn('a:p'))
    pPr_tmpl = None
    rPr_tmpl = None
    if first_p is not None:
        pPr = first_p.find(qn('a:pPr'))
        if pPr is not None:
            pPr_tmpl = copy.deepcopy(pPr)
        first_r = first_p.find(qn('a:r'))
        if first_r is not None:
            rPr = first_r.find(qn('a:rPr'))
            if rPr is not None:
                rPr_tmpl = copy.deepcopy(rPr)

    # Remove all existing paragraphs
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    # Add one paragraph per bullet
    for text in bullets:
        p_new = etree.SubElement(txBody, qn('a:p'))
        if pPr_tmpl is not None:
            p_new.append(copy.deepcopy(pPr_tmpl))
        r_new = etree.SubElement(p_new, qn('a:r'))
        if rPr_tmpl is not None:
            r_new.append(copy.deepcopy(rPr_tmpl))
        t_new = etree.SubElement(r_new, qn('a:t'))
        t_new.text = str(text)

    # Required trailing empty paragraph
    end_p = etree.SubElement(txBody, qn('a:p'))
    if pPr_tmpl is not None:
        end_p.append(copy.deepcopy(pPr_tmpl))


def _split_body(text):
    """Split body text into lines, handling |, ;, and newline delimiters."""
    for sep in ("|", ";", "\n"):
        text = text.replace(sep, "\n")
    return [l.strip() for l in text.split("\n") if l.strip()]


# ─── Slide operations ──────────────────────────────────────────────────────────

def _clone_slide(prs, src_idx):
    """
    Clone a slide and append to prs.
    Skips:
      - p:graphicFrame (charts — added fresh)
      - p:pic (picture shapes — have r:embed refs)
      - p:sp elements that contain r:embed/r:link refs (picture placeholders)
    Also copies the source slide's <p:bg> background override so cloned slides
    inherit the same background (e.g. white override on accent2-blue layout).
    Returns the new slide.
    """
    src       = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)

    dst_tree = new_slide.shapes._spTree
    src_tree = src.shapes._spTree

    # Clear default shapes added by add_slide
    for tag in (qn('p:sp'), qn('p:graphicFrame'), qn('p:pic'), qn('p:grpSp')):
        for elem in list(dst_tree.findall(tag)):
            dst_tree.remove(elem)

    # Copy elements from source, skipping anything with image references
    for elem in src_tree:
        tag = elem.tag
        if tag == qn('p:graphicFrame'):
            continue  # charts — we add fresh ones
        if tag == qn('p:pic'):
            continue  # standalone pictures — have r:embed
        if _has_image_ref(elem):
            continue  # picture placeholders with embedded images
        dst_tree.append(copy.deepcopy(elem))

    # Copy the source slide's background override (<p:bg>) so the layout's
    # background colour (accent2 blue on 1/3 Layout) doesn't bleed through.
    src_cSld = src._element.find(qn('p:cSld'))
    dst_cSld = new_slide._element.find(qn('p:cSld'))
    if src_cSld is not None and dst_cSld is not None:
        src_bg = src_cSld.find(qn('p:bg'))
        if src_bg is not None:
            # Remove any existing bg on the destination first
            dst_bg = dst_cSld.find(qn('p:bg'))
            if dst_bg is not None:
                dst_cSld.remove(dst_bg)
            # Insert bg as the first child of cSld (before spTree)
            dst_cSld.insert(0, copy.deepcopy(src_bg))

    return new_slide


def _delete_slide(prs, idx):
    """
    Remove slide at idx from the presentation.

    Cleans up:
    1. The slide entry in the presentation's sldIdLst
    2. The presentation → slide relationship
    3. The slide part itself from the package (prevents orphaned-part repair)
    4. Any chart parts (and their embedded xlsx) the slide owns exclusively
       (the template chart slide has a real embedded chart — deleting without
        cleaning its chart/xlsx parts triggers the PowerPoint repair dialog)

    Shared parts (slide layout, master, images) are intentionally left alone.
    """
    # Relationship types that are SHARED — never delete these targets
    _SHARED_RELTYPES = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
    }

    sldIdLst = prs.slides._sldIdLst
    sldId    = sldIdLst[idx]
    rId      = sldId.get(qn('r:id'))

    # Grab the slide part before severing the relationship
    slide_part = None
    try:
        slide_part = prs.part.related_parts.get(rId)
    except Exception:
        pass

    # Remove from slide id list
    sldIdLst.remove(sldId)

    # Remove presentation → slide relationship
    try:
        prs.part.drop_rel(rId)
    except AttributeError:
        try:
            del prs.part._rels[rId]
        except (KeyError, AttributeError, TypeError):
            pass

    if slide_part is None:
        return

    pkg = prs.part.package

    def _drop_from_pkg(part):
        """Remove a part from the package by matching partname string."""
        try:
            target_str = str(part.partname)
            for k in list(pkg._parts.keys()):
                if str(k) == target_str:
                    del pkg._parts[k]
                    break
        except Exception:
            pass

    # Collect slide-exclusive parts (charts, notes, etc.) — skip shared ones
    exclusive_parts = []
    try:
        for rel in slide_part._rels.values():
            if getattr(rel, 'is_external', False):
                continue
            if rel.reltype in _SHARED_RELTYPES:
                continue
            try:
                target = slide_part.related_parts.get(rel.rId)
                if target:
                    exclusive_parts.append(target)
            except Exception:
                pass
    except Exception:
        pass

    # Also collect sub-parts of exclusive parts (e.g. chart → embedded xlsx)
    sub_parts = []
    for ep in exclusive_parts:
        try:
            for sub_rel in ep._rels.values():
                if getattr(sub_rel, 'is_external', False):
                    continue
                try:
                    sub_target = ep.related_parts.get(sub_rel.rId)
                    if sub_target:
                        sub_parts.append(sub_target)
                except Exception:
                    pass
        except Exception:
            pass

    # Delete everything, deepest first
    for p in sub_parts:
        _drop_from_pkg(p)
    for p in exclusive_parts:
        _drop_from_pkg(p)
    _drop_from_pkg(slide_part)


def _suppress_layout_phs(slide, idxs):
    """
    Add empty slide-level placeholder overrides so layout instruction text
    (e.g. 'Chart goes here', 'Click to edit...') doesn't bleed through.
    """
    # Find which idxs are already overridden on the slide
    existing = {
        s.placeholder_format.idx
        for s in slide.shapes
        if s.is_placeholder
    }
    sp_tree = slide.shapes._spTree

    for idx in idxs:
        if idx in existing:
            continue
        sp_xml = (
            f'<p:sp'
            f' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{9000 + idx}" name="suppress_{idx}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
            f'<p:nvPr><p:ph idx="{idx}"/></p:nvPr>'
            f'</p:nvSpPr>'
            f'<p:spPr><a:noFill/></p:spPr>'
            f'<p:txBody>'
            f'<a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:endParaRPr lang="en-AU" dirty="0"/></a:p>'
            f'</p:txBody>'
            f'</p:sp>'
        )
        sp_tree.append(etree.fromstring(sp_xml))


# ─── Slide fill functions ──────────────────────────────────────────────────────

def _fill_cover(slide, campaign_name, date_range):
    _set_text(_ph(slide, 0),  campaign_name or "Post-Campaign Analysis")
    _set_text(_ph(slide, 15), date_range or "")


def _fill_agenda(slide, sections):
    for i, (title_idx, num_idx) in enumerate(AGENDA_ROWS):
        if i < len(sections):
            _set_text(_ph(slide, num_idx),   f"{i + 1:02d}")
            _set_text(_ph(slide, title_idx), sections[i])
        else:
            _set_text(_ph(slide, num_idx),   "")
            _set_text(_ph(slide, title_idx), "")


def _fill_divider(slide, section_title, section_num):
    _set_text(_ph(slide, DIVIDER_NUM_IDX),   f"{section_num:02d}")
    _set_text(_ph(slide, DIVIDER_TITLE_IDX), section_title)
    _set_text(_ph(slide, DIVIDER_SUB_IDX),   "")


def _fill_chart_slide(slide, slide_data):
    """
    Fills a cloned 1/3 Layout slide.
    - Title 1 (wide textbox) carries the headline — the visually dominant shape
    - Title 3 (narrow placeholder, ph_idx=0) is cleared to avoid double headers
    - Body bullets go into Text Placeholder 24
    - A fresh chart is added on the right
    - Inherited layout instruction placeholders are suppressed
    """
    title   = slide_data.get("title", "")
    bullets = list(slide_data.get("bullet_points", []) or [])
    if slide_data.get("body") and not bullets:
        bullets = [slide_data["body"]]
    so_what = slide_data.get("so_what", "")
    if so_what and so_what not in bullets:
        bullets.insert(0, f"💡 {so_what}")

    # Remove the narrow title placeholder (Title 3 / ph type="title") from the
    # spTree entirely — a visible empty placeholder conflicts with Title 1.
    # _suppress_layout_phs then adds a transparent no-fill override for ph=0
    # so the layout's "Title goes here" doesn't bleed through.
    sp_tree = slide.shapes._spTree
    for elem in list(sp_tree):
        if elem.tag == qn('p:sp'):
            ph = elem.find('.//' + qn('p:ph'))
            if ph is not None and ph.get('type') == 'title':
                sp_tree.remove(elem)
                break

    # Set headline in the full-width Title 1 textbox
    title_box = _shape(slide, "Title 1")
    if title_box and title_box.has_text_frame:
        _set_text(title_box, title)

    # Body bullets in Text Placeholder 24
    body = _shape(slide, "Text Placeholder 24")
    if body:
        _set_bullets(body, bullets)

    # Suppress layout instruction placeholders (transparent no-fill overrides)
    _suppress_layout_phs(slide, _CHART_SLIDE_SUPPRESS)

    # Add chart
    stype      = slide_data.get("slide_type", "data")
    chart_dict = slide_data.get("chart_data")
    chart_str  = (slide_data.get("chart_type") or "column").lower()
    xl_type    = CHART_TYPE_MAP.get(chart_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    if stype not in NO_CHART_SLIDE_TYPES and chart_dict:
        cd = CategoryChartData()
        cd.categories = chart_dict.get("categories", [])

        # Support both multi-series and legacy single-series formats
        if "series" in chart_dict and isinstance(chart_dict["series"], list):
            series_list = chart_dict["series"]
            for s in series_list:
                cd.add_series(s.get("name", "Metric"), _clean_values(s.get("values", [])))
            primary_label = series_list[0].get("name", "") if series_list else ""
        else:
            primary_label = chart_dict.get("series_name", "Metric")
            cd.add_series(primary_label, _clean_values(chart_dict.get("values", [])))

        gf = slide.shapes.add_chart(xl_type, CHART_L, CHART_T, CHART_W, CHART_H, cd)
        _apply_chart_style(gf.chart, xl_type, primary_label)


def _fill_rec_slide(slide, rec_data):
    """
    4-Col recommendations slide.
    Number slot → category heading ("What Worked", etc.)
    Body slot    → bullet points (4 points per column)
    """
    title_ph = slide.shapes.title
    if title_ph:
        _set_text(title_ph, "Recommendations")

    if not rec_data:
        return

    columns = rec_data.get("columns", [])
    for i, (num_idx, body_idx) in enumerate(REC_COLUMNS):
        col = columns[i] if i < len(columns) else {}
        # Put the category name in the number slot (replaces "01", "02"…)
        _set_text(_ph(slide, num_idx), col.get("title", f"{i + 1:02d}"))
        body_shape = _ph(slide, body_idx)
        if body_shape:
            lines = _split_body(col.get("body", ""))
            if lines:
                _set_bullets(body_shape, lines)


def _fill_callout(slide, callout_data):
    """Big Text callout slide — single bold statement."""
    if not callout_data:
        return
    statement = callout_data.get("statement") or callout_data.get("superhead", "")
    _set_text(_ph(slide, CALLOUT_STATEMENT_IDX), statement)


# ─── Section extraction ────────────────────────────────────────────────────────

def _extract_sections(slides_json):
    """
    Group slides into sections (max 5).
    Uses LLM-provided 'sections' + per-slide 'section' field if available.
    Falls back to auto-grouping.
    """
    sections_list = slides_json.get("sections", [])
    all_slides    = [s for s in slides_json.get("slides", [])
                     if s.get("slide_type") not in ("recommendation",)]

    if sections_list and all_slides:
        groups    = {s: [] for s in sections_list}
        ungrouped = []
        for sl in all_slides:
            sec = sl.get("section", "")
            if sec in groups:
                groups[sec].append(sl)
            else:
                ungrouped.append(sl)
        result = [{"title": t, "slides": groups[t]} for t in sections_list]
        if ungrouped and result:
            result[-1]["slides"].extend(ungrouped)
        return result[:5]

    # Auto-group fallback
    intro_slides = [s for s in all_slides if s.get("slide_type") == "intro"]
    data_slides  = [s for s in all_slides if s.get("slide_type") != "intro"]
    sections     = []
    if intro_slides:
        sections.append({"title": "Overview", "slides": intro_slides})
    n_chunks = min(4, max(1, len(data_slides) // 4))
    chunk    = max(1, len(data_slides) // n_chunks)
    for i in range(0, len(data_slides), chunk):
        batch = data_slides[i:i + chunk]
        if batch:
            sections.append({"title": f"Performance Analysis {len(sections)}", "slides": batch})
        if len(sections) >= 5:
            break
    return sections[:5]


# ─── Slide renumbering ─────────────────────────────────────────────────────────

def _renumber_slides(path):
    """
    Post-process the saved PPTX zip so slide XML files are numbered
    consecutively (slide1.xml, slide2.xml, …) with no gaps.

    Background: the template has slides 1–6.  We add content slides which get
    numbers 7, 8, 9 … then delete template slides 3–6, leaving gaps like
    1, 2, 7, 8 … PowerPoint treats non-consecutive slide numbers as a corrupt
    package and prompts for repair.  This step rewrites the zip to fix it.

    Uses a single-pass regex substitution to avoid cascading replacement bugs
    (e.g. slide9→slide5 being re-matched by the slide5 rule on the next pass).
    """
    _SLIDE_RE = re.compile(r'^ppt/slides/slide(\d+)\.xml$')

    # ── Read every file out of the zip ───────────────────────────────────────
    with zipfile.ZipFile(path, 'r') as zin:
        file_data = {name: zin.read(name) for name in zin.namelist()}

    # ── Find slide numbers present ────────────────────────────────────────────
    old_nums = sorted(
        int(m.group(1))
        for name in file_data
        for m in [_SLIDE_RE.match(name)] if m
    )

    if old_nums == list(range(1, len(old_nums) + 1)):
        return  # already consecutive — nothing to do

    # ── Build old→new number map ──────────────────────────────────────────────
    num_map = {old: new for new, old in enumerate(old_nums, 1) if old != new}
    if not num_map:
        return

    # ── Build filename rename map ─────────────────────────────────────────────
    rename = {}
    for old_num, new_num in num_map.items():
        for tmpl in ('ppt/slides/slide{}.xml',
                     'ppt/slides/_rels/slide{}.xml.rels'):
            old_p = tmpl.format(old_num)
            new_p = tmpl.format(new_num)
            if old_p in file_data:
                rename[old_p] = new_p

    # ── Single-pass regex substitution for text files ────────────────────────
    # Matches "slideN.xml" or "slideN.xml.rels" where N is any old slide number.
    # A single-pass approach avoids the conflict where slide9→slide5 and then
    # the literal text "slide5" gets matched by the slide5→slide2 rule.
    _pat = re.compile(
        r'slide(' + '|'.join(str(n) for n in sorted(num_map, reverse=True)) + r')(\.xml(?:\.rels)?)'
    )
    def _replacer(m):
        old_n = int(m.group(1))
        return f'slide{num_map[old_n]}{m.group(2)}'

    def _apply(text):
        return _pat.sub(_replacer, text)

    # ── Rewrite zip ───────────────────────────────────────────────────────────
    TEXT_SUFFIXES = ('.xml', '.rels')
    new_data = {}
    for name, content in file_data.items():
        new_name = rename.get(name, name)
        if any(name.endswith(s) for s in TEXT_SUFFIXES):
            content = _apply(content.decode('utf-8')).encode('utf-8')
        new_data[new_name] = content

    tmp = path + '.renumber_tmp'
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, content in new_data.items():
            zout.writestr(name, content)
    os.replace(tmp, path)


# ─── Main builder ──────────────────────────────────────────────────────────────

def build_pptx_from_template(
    slides_json: dict,
    output_path: str,
    campaign_name: str = "",
    date_range: str = ""
) -> str:
    """
    Build a branded PPTX from OMD template.pptx.

    Deck order: Cover → Agenda → [Divider + data slides] × N → Rec → Callout
    """
    prs = Presentation(TEMPLATE_PATH)

    # ── Cover ────────────────────────────────────────────────
    _fill_cover(prs.slides[TMPL_COVER], campaign_name, date_range)

    # ── Agenda ───────────────────────────────────────────────
    sections = _extract_sections(slides_json)
    _fill_agenda(prs.slides[TMPL_AGENDA], [s["title"] for s in sections])

    # ── Sections ─────────────────────────────────────────────
    for sec_idx, section in enumerate(sections):
        div = _clone_slide(prs, TMPL_DIVIDER)
        _fill_divider(div, section["title"], sec_idx + 1)

        for slide_data in section["slides"]:
            cs = _clone_slide(prs, TMPL_CHART)
            _fill_chart_slide(cs, slide_data)

    # ── Recommendations — 4-col layout ───────────────────────
    rec = _clone_slide(prs, TMPL_REC)
    _fill_rec_slide(rec, slides_json.get("recommendations"))

    # ── Callout ──────────────────────────────────────────────
    co = _clone_slide(prs, TMPL_CALLOUT)
    _fill_callout(co, slides_json.get("callout"))

    # ── Remove original template source slides (reverse order) ──
    for idx in (TMPL_CALLOUT, TMPL_REC, TMPL_CHART, TMPL_DIVIDER):
        _delete_slide(prs, idx)

    prs.save(output_path)
    _renumber_slides(output_path)   # fix non-consecutive slide numbers → no repair dialog
    return output_path


# ─── Backwards compat ─────────────────────────────────────────────────────────

def build_pptx_from_scratch(slides_json: dict, output_path: str) -> str:
    return build_pptx_from_template(slides_json, output_path)


# ─── Slide Generator builder ───────────────────────────────────────────────────

_SG_TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "slide generator.pptx")

def build_quick_slides_pptx(slides_json: dict, output_path: str) -> str:
    """
    Build a PPTX from the 'slide generator.pptx' template.

    The template has a single 1/3 Layout slide (chart + left commentary).
    We clone it once per slide in slides_json["slides"], fill content,
    then delete the original template slide and renumber.

    No cover / agenda / dividers / rec — just the content slides.
    """
    prs = Presentation(_SG_TEMPLATE_PATH)

    TMPL = 0   # the single template slide

    for slide_data in slides_json.get("slides", []):
        cs = _clone_slide(prs, TMPL)
        _fill_chart_slide(cs, slide_data)

    # Delete the original template slide
    if prs.slides:
        _delete_slide(prs, TMPL)

    prs.save(output_path)
    _renumber_slides(output_path)
    return output_path


# ─── Media Strategy PPTX builder ──────────────────────────────────────────────

def build_media_strategy_pptx(strategy: dict, output_path: str) -> str:
    """
    Build a branded strategy deck from a media strategy JSON object.
    Uses the OMD template (cover + divider + chart slides).
    Produces: Cover → Strategy Overview → Channel Mix → Budget → Flight Plan → Insights.
    """
    prs = Presentation(TEMPLATE_PATH)

    client_name = strategy.get("client_name", "Client")
    headline    = strategy.get("strategy_headline", "Media Strategy")
    bs          = strategy.get("brief_summary", {})
    sa          = strategy.get("strategic_approach", {})

    # ── Cover ──────────────────────────────────────────────────
    cover = _clone_slide(prs, TMPL_COVER)
    for ph in cover.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            _set_text(ph, f"Media Strategy — {client_name}")
        elif idx == 1:
            _set_text(ph, headline)

    # ── Agenda ─────────────────────────────────────────────────
    agenda = _clone_slide(prs, TMPL_AGENDA)
    sections = [
        ("01", "Strategy Overview"),
        ("02", "Channel Mix & Rationale"),
        ("03", "Budget Allocation"),
        ("04", "Monthly Flight Plan"),
        ("05", "Insights & Recommendations"),
    ]
    for (num_idx, title_idx), (num, title) in zip(AGENDA_ROWS, sections):
        n = _ph(agenda, num_idx)
        t = _ph(agenda, title_idx)
        if n: _set_text(n, num)
        if t: _set_text(t, title)

    # ─── Helper: strategy divider slide ───────────────────────
    def _divider(num, title, sub=""):
        sl = _clone_slide(prs, TMPL_DIVIDER)
        n  = _ph(sl, DIVIDER_NUM_IDX)
        t  = _ph(sl, DIVIDER_TITLE_IDX)
        s  = _ph(sl, DIVIDER_SUB_IDX)
        if n: _set_text(n, num)
        if t: _set_text(t, title)
        if s: _set_text(s, sub)

    # ─── Helper: bullet data slide ────────────────────────────
    def _bullet_slide(title, bullets, chart_type=None, chart_data=None):
        sl_data = {
            "title":        title,
            "slide_type":   "data" if chart_data else "intro",
            "bullet_points": bullets,
            "so_what":      "",
            "chart_type":   chart_type,
            "chart_data":   chart_data,
        }
        sl = _clone_slide(prs, TMPL_CHART)
        _fill_chart_slide(sl, sl_data)

    # ══ SECTION 1: Strategy Overview ════════════════════════════
    _divider("01", "Strategy Overview", bs.get("campaign_period", ""))
    _bullet_slide(
        "Campaign Strategy",
        [
            f"Objective: {obj}" for obj in bs.get("objectives", [])
        ] + [
            f"Target Audience: {bs.get('target_audience', '')}",
            f"Positioning: {sa.get('positioning', '')}",
            f"Channel Philosophy: {sa.get('channel_philosophy', '')}",
        ]
    )
    _bullet_slide(
        "Strategic Approach",
        [strategy.get("executive_summary", "")]
        + [f"Key Tension: {t}" for t in sa.get("key_tensions", [])]
        + [f"Key Message: {m}" for m in bs.get("key_messages", [])],
    )

    # ══ SECTION 2: Channel Mix ════════════════════════════════════
    _divider("02", "Channel Mix & Rationale")

    channel_mix = strategy.get("channel_mix", [])
    if channel_mix:
        # Overview chart: budget % by channel
        ch_categories  = [c.get("channel", "") for c in channel_mix]
        ch_values      = [c.get("budget_pct", 0) for c in channel_mix]
        _bullet_slide(
            "Recommended Channel Mix",
            [f"{c.get('channel')}: {c.get('budget_pct')}% — {c.get('role', '')}" for c in channel_mix],
            chart_type="bar",
            chart_data={"categories": ch_categories, "values": ch_values, "series_name": "Budget %"},
        )
        # Individual channel slides
        for ch in channel_mix:
            bullets = [
                f"Role: {ch.get('role', '')}",
                f"Rationale: {ch.get('rationale', '')}",
                f"Platforms: {', '.join(ch.get('platforms', []))}",
                f"Formats: {', '.join(ch.get('recommended_formats', []))}",
                f"Past Performance: {ch.get('past_performance', '')}",
                f"Market Benchmark: {ch.get('market_benchmark', '')}",
            ]
            _bullet_slide(ch.get("channel", "Channel"), [b for b in bullets if b.split(": ", 1)[-1].strip()])

    # ══ SECTION 3: Budget Allocation ════════════════════════════
    _divider("03", "Budget Allocation")
    alloc = strategy.get("budget_allocation", [])
    if alloc:
        _bullet_slide(
            "Budget Allocation by Channel",
            [f"{a.get('channel')}: {a.get('pct')}%  ({a.get('est_spend', '')})" for a in alloc],
            chart_type="bar",
            chart_data={
                "categories": [a.get("channel", "") for a in alloc],
                "values":     [a.get("pct", 0) for a in alloc],
                "series_name": "Budget %",
            },
        )

    # ══ SECTION 4: Monthly Flight Plan ════════════════════════════
    _divider("04", "Monthly Flight Plan", bs.get("campaign_period", ""))
    flight = strategy.get("monthly_flight", [])
    if flight:
        # Split into two halves if more than 6 months
        half = (len(flight) + 1) // 2
        for chunk_start in range(0, len(flight), half):
            chunk = flight[chunk_start:chunk_start + half]
            months = [r.get("month", "") for r in chunk]
            weights = {"Heavy": 3, "Medium": 2, "Light": 1, "Off": 0}
            values  = [weights.get(r.get("relative_weight", "Medium"), 1) for r in chunk]
            bullets = [f"{r.get('month')}: {r.get('phase', '')} — {r.get('activity', '')}" for r in chunk]
            _bullet_slide(
                f"Flight Plan — {months[0]} to {months[-1]}",
                bullets,
                chart_type="column",
                chart_data={"categories": months, "values": values, "series_name": "Activity Weight (Heavy=3)"},
            )

    # ══ SECTION 5: Insights & Recommendations ════════════════════
    _divider("05", "Insights & Recommendations")

    market_insights = strategy.get("market_insights", [])
    if market_insights:
        _bullet_slide(
            "Market Insights from Portfolio",
            [f"{i.get('insight', '')}  →  {i.get('implication', '')}" for i in market_insights],
        )

    learnings = strategy.get("past_campaign_learnings", [])
    if learnings:
        _bullet_slide(
            "Learnings from Past Campaigns",
            [f"{l.get('learning', '')}  ➔  {l.get('applied_as', '')}" for l in learnings],
        )

    recs = strategy.get("recommendations", [])
    if recs:
        _bullet_slide(
            "Recommendations",
            [f"[{r.get('priority', '')}] {r.get('recommendation', '')} — {r.get('rationale', '')}" for r in recs],
        )

    risks = strategy.get("risks_and_mitigations", [])
    if risks:
        _bullet_slide(
            "Risks & Mitigations",
            [f"Risk: {r.get('risk', '')}  |  Mitigation: {r.get('mitigation', '')}" for r in risks],
        )

    # ── Rec summary slide ───────────────────────────────────────
    rec_slide = _clone_slide(prs, TMPL_REC)
    rec_cols = [
        {"number": "01", "title": "Objectives",        "body": "\n".join(bs.get("objectives", []))},
        {"number": "02", "title": "Audience",           "body": bs.get("target_audience", "")},
        {"number": "03", "title": "Channel Philosophy", "body": sa.get("channel_philosophy", "")},
        {"number": "04", "title": "Top Recommendation", "body": recs[0].get("recommendation", "") if recs else ""},
    ]
    for (num_idx, body_idx), col_data in zip(REC_COLUMNS, rec_cols):
        n = _ph(rec_slide, num_idx)
        b = _ph(rec_slide, body_idx)
        if n: _set_text(n, col_data["number"])
        if b: _set_text(b, col_data["body"])

    # ── Callout ─────────────────────────────────────────────────
    callout_data = {
        "superhead": "The Strategy in One Line",
        "statement": headline,
    }
    co = _clone_slide(prs, TMPL_CALLOUT)
    _fill_callout(co, callout_data)

    # ── Remove template source slides ──────────────────────────
    for idx in (TMPL_CALLOUT, TMPL_REC, TMPL_CHART, TMPL_DIVIDER, TMPL_AGENDA, TMPL_COVER):
        _delete_slide(prs, idx)

    prs.save(output_path)
    _renumber_slides(output_path)
    return output_path
